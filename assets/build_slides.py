"""
Build the Atlanta Robotics PowerPoint presentation.
Run: python3 assets/build_slides.py
Output: assets/Atlanta_Robotics_DataCenter_DigitalTwin.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, os

OUT = os.path.join(os.path.dirname(__file__), "Atlanta_Robotics_DataCenter_DigitalTwin.pptx")

# ── Brand colours ──────────────────────────────────────────────────────────────
NVIDIA_GREEN  = RGBColor(0x76, 0xB9, 0x00)   # #76B900
DARK_BG       = RGBColor(0x0F, 0x0F, 0x0F)   # near-black
SLIDE_BG      = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT        = RGBColor(0x00, 0xD4, 0xFF)   # cyan accent
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY    = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE        = RGBColor(0xFF, 0x8C, 0x00)
CARD_BG       = RGBColor(0x16, 0x21, 0x3E)   # slightly lighter navy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helper functions ───────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=SLIDE_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return box

def heading(slide, text, top=0.35, size=36, color=NVIDIA_GREEN):
    txbox(slide, text, 0.5, top, 12.33, 0.7, size=size, bold=True, color=color, align=PP_ALIGN.LEFT)

def subheading(slide, text, top=0.9, size=22, color=ACCENT):
    txbox(slide, text, 0.5, top, 12.33, 0.5, size=size, bold=False, color=color, align=PP_ALIGN.LEFT)

def divider(slide, top=1.15, color=NVIDIA_GREEN):
    rect(slide, 0.5, top, 12.33, 0.04, fill_color=color)

def bullet_block(slide, items, l, t, w, h, size=16, title=None, title_color=ACCENT):
    if title:
        txbox(slide, title, l, t, w, 0.35, size=17, bold=True, color=title_color)
        t += 0.38
        h -= 0.38
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(size)
        run.font.color.rgb = LIGHT_GREY
    return box

def card(slide, l, t, w, h, title, body_lines, title_color=ACCENT, body_size=14):
    rect(slide, l, t, w, h, fill_color=CARD_BG, line_color=ACCENT, line_width=Pt(0.75))
    txbox(slide, title, l+0.12, t+0.1, w-0.24, 0.38, size=15, bold=True, color=title_color, align=PP_ALIGN.LEFT)
    rect(slide, l+0.12, t+0.48, w-0.24, 0.03, fill_color=ACCENT)
    bullet_block(slide, body_lines, l+0.12, t+0.55, w-0.24, h-0.65, size=body_size)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)

# Top accent bar
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)

# Large title
txbox(s, "Data Center Digital Twin", 0.6, 1.2, 12, 1.1,
      size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txbox(s, "→  AI World Model on Google Cloud", 0.6, 2.25, 12, 0.7,
      size=34, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)

# Horizontal rule
rect(s, 0.6, 3.05, 11.5, 0.05, fill_color=NVIDIA_GREEN)

# Subtitle line
txbox(s, "Synthetic failure data  ·  Temporal Transformer  ·  Vertex AI  ·  Omniverse Streaming",
      0.6, 3.2, 12, 0.5, size=18, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

# Bottom tag
txbox(s, "Atlanta Robotics  —  College Student Tutorial Series",
      0.6, 6.5, 8, 0.45, size=16, color=NVIDIA_GREEN, align=PP_ALIGN.LEFT)
txbox(s, "2026", 11.5, 6.5, 1.5, 0.45, size=16, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)

rect(s, 0, 7.42, 13.33, 0.08, fill_color=NVIDIA_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "What We Will Cover Today")
divider(s)

phases = [
    ("01", "What is a Digital Twin?",          "The problem, the solution, USD format"),
    ("02", "System Architecture",              "How every piece connects end-to-end"),
    ("03", "GCP Infrastructure",               "Bucket · Artifact Registry · GPU VM · Vertex AI"),
    ("04", "Synthetic Data Generation",        "Simulating overheating, disk failure, cooling loss"),
    ("05", "Temporal Transformer World Model", "Attention · per-horizon heads · training"),
    ("06", "Live Inference Results",           "What the model actually predicts"),
    ("07", "Student Workflow",                 "One command to get started"),
    ("08", "Exercises & Next Steps",           "What you build from here"),
]

col_w, col_h = 5.7, 0.72
positions = [
    (0.4,  1.3), (6.9,  1.3),
    (0.4,  2.1), (6.9,  2.1),
    (0.4,  2.9), (6.9,  2.9),
    (0.4,  3.7), (6.9,  3.7),
]

for (num, title, subtitle), (cx, cy) in zip(phases, positions):
    rect(s, cx, cy, col_w, col_h-0.06, fill_color=CARD_BG, line_color=ACCENT, line_width=Pt(0.5))
    txbox(s, num, cx+0.1, cy+0.08, 0.5, 0.35, size=18, bold=True, color=NVIDIA_GREEN)
    txbox(s, title, cx+0.58, cy+0.06, col_w-0.7, 0.32, size=14, bold=True, color=WHITE)
    txbox(s, subtitle, cx+0.58, cy+0.36, col_w-0.7, 0.28, size=11, color=LIGHT_GREY, italic=True)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial",
      0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "2 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "The Problem We Are Solving")
subheading(s, "Why do we need a digital twin + AI? Why not just monitor sensors?")
divider(s)

# Left column — problem
rect(s, 0.4, 1.3, 5.8, 5.3, fill_color=CARD_BG, line_color=RGBColor(0xCC,0x44,0x44), line_width=Pt(1))
txbox(s, "❌  The Real-World Data Problem", 0.55, 1.4, 5.5, 0.4, size=15, bold=True, color=RGBColor(0xFF,0x66,0x66))
rect(s, 0.55, 1.8, 5.5, 0.03, fill_color=RGBColor(0xCC,0x44,0x44))

problems = [
    "5,000 servers — maybe 10 failures per year",
    "That is only 0.01% failure rate in training data",
    "Model trained on this just predicts 'normal' always",
    "99.99% accuracy — completely useless for alerting",
    "",
    "Real failures are dangerous to allow on purpose",
    "No company will let you stress-test live hardware",
    "Historical logs are confidential and hard to get",
    "",
    "Result: you cannot train a reliable failure model",
    "on real data alone — there is not enough of it",
]
bullet_block(s, problems, 0.55, 1.9, 5.5, 4.5, size=14)

# Right column — solution
rect(s, 6.9, 1.3, 5.9, 5.3, fill_color=CARD_BG, line_color=NVIDIA_GREEN, line_width=Pt(1))
txbox(s, "✅  The Digital Twin Solution", 7.05, 1.4, 5.6, 0.4, size=15, bold=True, color=NVIDIA_GREEN)
rect(s, 7.05, 1.8, 5.6, 0.03, fill_color=NVIDIA_GREEN)

solutions = [
    "Build a virtual replica of the data center",
    "Simulate failures as many times as you want",
    "Generate thousands of labeled failure events",
    "No real hardware at risk — all virtual",
    "",
    "Model trains on synthetic data",
    "Deploys against real sensor streams",
    "Works because simulation matches physics",
    "",
    "Used by Siemens, GE, NVIDIA in production",
    "This is the exact same workflow",
]
bullet_block(s, solutions, 7.05, 1.9, 5.6, 4.5, size=14)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "3 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What is USD
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "USD — Universal Scene Description")
subheading(s, "The 3D file format that powers the digital twin  (created by Pixar, adopted by NVIDIA)")
divider(s)

txbox(s, "Think of USD like a web page with links — the root file is tiny but references thousands of sub-assets.",
      0.5, 1.25, 12.3, 0.5, size=16, color=LIGHT_GREY)

# Three concept cards
for i, (title, lines) in enumerate([
    ("Composition", [
        "A scene is assembled from many files",
        "DataHall_Full_01.usd = 31 KB",
        "But it references 1,822 sub-files",
        "Total asset tree = 9.6 GB",
        "Root file just says: place this rack here",
    ]),
    ("Instancing", [
        "One DGX A100 model = ~800 MB",
        "Referenced 48 times, not copied",
        "GPU loads geometry once",
        "Then transforms it 48 times",
        "Saves gigabytes of memory",
    ]),
    ("Layers", [
        "Multiple USD files can override each other",
        "Base layer: geometry + materials",
        "Override layer: sensor attribute values",
        "Animation layer: failure simulation",
        "Non-destructive — originals unchanged",
    ]),
]):
    card(s, 0.4 + i*4.3, 1.85, 4.1, 3.0, title, lines, body_size=13)

# Asset tree diagram
rect(s, 0.4, 5.0, 12.5, 2.15, fill_color=RGBColor(0x0D, 0x1B, 0x2A), line_color=ACCENT, line_width=Pt(0.5))
tree = (
    "Datacenter_NVD@10012/                    ← 9.6 GB, 1,822 files\n"
    "  └── Assets/DigitalTwin/Assets/Datacenter/\n"
    "        ├── Racks/Rack_42U_A/            ← 42U rack shell + materials\n"
    "        ├── DGX_Servers/A100/            ← NVIDIA DGX A100 3D model + 4K textures\n"
    "        ├── Network_Switches/QM8700/     ← InfiniBand switch models\n"
    "        ├── Power_Distribution/PDU_A/    ← Power distribution units\n"
    "        └── Facilities/Stages/Data_Hall/\n"
    "              └── DataHall_Full_01.usd   ← ROOT FILE (31 KB — just references above)"
)
txbox(s, tree, 0.55, 5.1, 12.2, 2.0, size=11, color=NVIDIA_GREEN)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "4 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "System Architecture — End to End")
divider(s, top=0.95)

arch = """\
 LOCAL MACHINE                     GOOGLE CLOUD PLATFORM
 ┌──────────────────────┐          ┌──────────────────────────────────────────────┐
 │  kit-app-template    │          │                                              │
 │  ├── USD Explorer    │─────────►│  Artifact Registry  (Docker images)          │
 │  └── DataHall config │          │                                              │
 │                      │          │  GCS Bucket  gs://project-omniverse-assets/  │
 │  DataHall_Full_01.usd│─────────►│  ├── Datacenter_NVD@10012/  (9.6 GB)        │
 │  (9.6 GB on disk)    │          │  ├── training-data/sensor_timeseries.csv     │
 └──────────────────────┘          │  └── models/best_model.pt                   │
                                   │                                              │
 STUDENT MACHINE                   │  GPU VM  (g2-standard-8 + NVIDIA L4)        │
 ┌──────────────────────┐          │  └── Kit App Container ──► Port 8011 (web)  │
 │  student_setup.sh    │─────────►│                                              │
 │  → downloads assets  │          │  Vertex AI                                   │
 │  → trains model      │          │  ├── Training Job  (A100)                    │
 └──────────────────────┘          │  └── Endpoint  (live failure probability)    │
                                   └──────────────────────────────────────────────┘
                                                    │
                                    failure probability JSON
                                                    │
                                    ┌───────────────▼─────────────┐
                                    │  Browser  http://VM-IP:8011 │
                                    │  Failing racks glow RED     │
                                    └─────────────────────────────┘"""

rect(s, 0.3, 1.05, 12.73, 6.1, fill_color=RGBColor(0x0D,0x1B,0x2A), line_color=ACCENT, line_width=Pt(0.5))
txbox(s, arch, 0.45, 1.1, 12.5, 6.0, size=11, color=NVIDIA_GREEN)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "5 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — GCP Infrastructure
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "Google Cloud Infrastructure — Phase 2")
subheading(s, "Four services, one script: bash deploy/02_gcp_setup.sh")
divider(s)

services = [
    ("Cloud Storage (GCS)",
     ["Object storage — like a cloud hard drive",
      "Stores the 9.6 GB USD stage",
      "Stores training CSV + trained model",
      "Students download from here",
      "Cost: ~$0.20 / month for 10 GB"]),
    ("Artifact Registry",
     ["Private Docker image registry",
      "Stores the built Kit container",
      "Same region as VM = fast pull",
      "No egress costs within region",
      "Like DockerHub but private + fast"]),
    ("Compute Engine VM",
     ["g2-standard-8 machine type",
      "8 vCPU, 32 GB RAM",
      "1× NVIDIA L4 GPU (24 GB VRAM)",
      "Runs the Kit streaming container",
      "Cost: ~$0.40 / hour"]),
    ("Vertex AI",
     ["Managed ML training platform",
      "Submit job → A100 auto-provisioned",
      "Train → save model → VM shuts down",
      "Endpoint serves live predictions",
      "Cost: ~$6 for 2-hour A100 job"]),
]
for i, (title, lines) in enumerate(services):
    col = i % 2
    row = i // 2
    card(s, 0.4 + col*6.4, 1.3 + row*2.85, 6.1, 2.7, title, lines, body_size=14)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "6 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Synthetic Data Generation
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "Synthetic Data Generation — Phase 6")
subheading(s, "48 racks × 30 days × 4 failure types  →  414,720 labeled sensor readings")
divider(s)

failures = [
    ("🌡  Overheating",
     ["temp_c: +25 to +45°C above baseline",
      "power_kw: +0.5 to +1.5 kW",
      "cpu_load: +10 to +30% (throttling)",
      "Cause: CRAC unit failure",
      "Onset: rapid (minutes to hours)"]),
    ("💾  Disk Degradation",
     ["disk_health: drops -40 to -70%",
      "temp_c: slight +2 to +8°C rise",
      "cpu_load: +5% (I/O wait)",
      "Cause: SMART failures, wear",
      "Onset: slow (weeks to months)"]),
    ("⚡  Power Fluctuation",
     ["power_kw: spike +2 to +5 kW",
      "temp_c: brief +3 to +10°C",
      "Other sensors: minimal change",
      "Cause: PDU failure, UPS switchover",
      "Onset: fast (seconds)"]),
    ("❄  Cooling Failure",
     ["temp_c: +15 to +35°C rise",
      "power_kw: +1 to +2.5 kW (fans max)",
      "disk_health: -10 to -30% (heat stress)",
      "Cause: chiller failure, coolant leak",
      "Onset: medium (hours)"]),
]
for i, (title, lines) in enumerate(failures):
    col = i % 2
    row = i // 2
    card(s, 0.4 + col*6.4, 1.3 + row*2.45, 6.1, 2.3, title, lines, body_size=13)

# Bell-curve ramp note
rect(s, 0.4, 6.25, 12.5, 0.75, fill_color=RGBColor(0x0D,0x1B,0x2A), line_color=ACCENT, line_width=Pt(0.5))
txbox(s, "Bell-curve ramping:  failures don't appear instantly — sensors ramp up then back down using sin(π·progress), matching real hardware behaviour",
      0.6, 6.32, 12.1, 0.62, size=13, color=LIGHT_GREY, italic=True)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "7 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Dataset Statistics
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "Dataset Statistics")
subheading(s, "What the generated CSV looks like")
divider(s)

# CSV sample
rect(s, 0.4, 1.25, 12.5, 1.7, fill_color=RGBColor(0x0D,0x1B,0x2A), line_color=ACCENT, line_width=Pt(0.5))
txbox(s, "Sample rows from sensor_timeseries.csv", 0.6, 1.3, 7, 0.3, size=12, color=ACCENT, bold=True)
csv_sample = (
    "timestamp               rack_id  temp_c  power_kw  disk_health  cpu_load  label\n"
    "2026-01-01T00:00:00          0    23.4      5.82       0.970      0.420    normal\n"
    "2026-01-15T08:20:00         12    48.7      7.31       0.910      0.730    overheating\n"
    "2026-02-03T14:45:00         31    25.1      5.63       0.412      0.581    disk_degradation\n"
    "2026-01-22T02:10:00          7    29.4     11.20       0.930      0.410    power_fluctuation"
)
txbox(s, csv_sample, 0.55, 1.6, 12.2, 1.3, size=11, color=NVIDIA_GREEN)

# Stats cards
stats = [
    ("Total Rows",    "414,720",   "48 racks × 8,640 steps"),
    ("Time Span",     "30 days",   "5-minute intervals"),
    ("Normal",        "96.9%",     "401,976 rows"),
    ("Failure",       "3.1%",      "12,744 rows (4 types)"),
    ("Training Set",  "360,288",   "90% of sliding windows"),
    ("Validation",    "40,032",    "10% held out"),
]
for i, (label, value, note) in enumerate(stats):
    col = i % 3
    row = i // 3
    rect(s, 0.4 + col*4.3, 3.1 + row*1.7, 4.0, 1.55, fill_color=CARD_BG, line_color=ACCENT, line_width=Pt(0.5))
    txbox(s, label, 0.6 + col*4.3, 3.18, 3.7, 0.3, size=12, color=LIGHT_GREY)
    txbox(s, value, 0.6 + col*4.3, 3.48, 3.7, 0.55, size=28, bold=True, color=NVIDIA_GREEN)
    txbox(s, note,  0.6 + col*4.3, 3.98, 3.7, 0.3,  size=11, color=LIGHT_GREY, italic=True)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "8 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Transformer Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "Temporal Transformer World Model — Phase 7")
subheading(s, "Why a Transformer beats an LSTM for this task")
divider(s)

# Architecture diagram (text-art)
rect(s, 0.4, 1.25, 7.0, 5.9, fill_color=RGBColor(0x0D,0x1B,0x2A), line_color=ACCENT, line_width=Pt(0.5))
arch_txt = """\
Input:  (batch=256, window=12, features=4)
  temp_c  power_kw  disk_health  cpu_load
        │
        ▼  Linear(4 → 64)
  (256, 12, 64)  ← project to d_model
        │
        ▼  Positional Encoding
  adds time-order signal to each step
        │
        ▼  Transformer Encoder  ×3 layers
        │   MultiHeadAttention (4 heads)
        │   each head learns different pattern:
        │   head 1 → temperature trends
        │   head 2 → power spikes
        │   head 3 → disk + temp correlation
        │   head 4 → cooling oscillations
        │
        ▼  Mean Pooling over 12 time steps
  (256, 64)  ← single vector per rack
        │
   ┌────┴────┬────────┐
   ▼         ▼        ▼
  MLP       MLP      MLP
   │         │        │
  1h        6h       24h
failure  failure  failure
  prob     prob     prob"""
txbox(s, arch_txt, 0.55, 1.32, 6.7, 5.75, size=11, color=NVIDIA_GREEN)

# Right column — key concepts
for i, (title, lines) in enumerate([
    ("Why Transformer, not LSTM?", [
        "LSTM reads sequence step-by-step",
        "Forgets events far in the past",
        "Cannot parallelize (slow to train)",
        "Transformer: every step attends to",
        "every other step simultaneously",
        "Faster + better long-range memory",
    ]),
    ("Three Prediction Horizons", [
        "1h  — imminent failures (3.6% positive)",
        "6h  — near-term    (6.7% positive)",
        "24h — early warning (17.3% positive)",
        "Separate MLP head per horizon",
        "Shared encoder = shared learning",
        "Per-horizon class weights applied",
    ]),
]):
    card(s, 7.6, 1.25 + i*3.0, 5.3, 2.8, title, lines, body_size=13)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "9 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Training Results
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "Training Results")
subheading(s, "RTX 3070 Laptop · 15 epochs · ~3 minutes total · batch size 512")
divider(s)

# Epoch log
rect(s, 0.4, 1.25, 12.5, 2.3, fill_color=RGBColor(0x0D,0x1B,0x2A), line_color=ACCENT, line_width=Pt(0.5))
txbox(s, "Training log (abridged)", 0.6, 1.3, 5, 0.3, size=12, color=ACCENT, bold=True)
log = (
    "Epoch   1/15  train_loss=1.9733  val_loss=0.6198  acc: 1h=90.0%  6h=87.8%  24h=80.3%  ← Saved\n"
    "Epoch   5/15  train_loss=1.8197  val_loss=0.5915  acc: 1h=86.0%  6h=84.4%  24h=75.0%  ← Saved\n"
    "Epoch  10/15  train_loss=1.8012  val_loss=0.5858  acc: 1h=88.4%  6h=85.9%  24h=76.6%\n"
    "Epoch  15/15  train_loss=1.7980  val_loss=0.5844  acc: 1h=89.0%  6h=86.6%  24h=77.1%  ← Best"
)
txbox(s, log, 0.55, 1.6, 12.2, 1.85, size=12, color=NVIDIA_GREEN)

# Accuracy cards
metrics = [
    ("1h Accuracy",  "89.0%", "Imminent failure detection"),
    ("6h Accuracy",  "86.6%", "Near-term prediction"),
    ("24h Accuracy", "77.1%", "Early warning (hardest)"),
]
for i, (label, value, note) in enumerate(metrics):
    rect(s, 0.4 + i*4.3, 3.7, 4.0, 1.6, fill_color=CARD_BG, line_color=NVIDIA_GREEN, line_width=Pt(1))
    txbox(s, label, 0.6 + i*4.3, 3.78, 3.7, 0.32, size=14, color=LIGHT_GREY)
    txbox(s, value, 0.6 + i*4.3, 4.08, 3.7, 0.65, size=36, bold=True, color=NVIDIA_GREEN)
    txbox(s, note,  0.6 + i*4.3, 4.7,  3.7, 0.32, size=11, color=LIGHT_GREY, italic=True)

# Note about accuracy
txbox(s, "Note: 1h > 6h > 24h accuracy is expected — predicting closer events is always easier.",
      0.5, 5.42, 12.3, 0.38, size=13, color=LIGHT_GREY, italic=True)

# Class weights note
rect(s, 0.4, 5.85, 12.5, 0.8, fill_color=CARD_BG, line_color=ACCENT, line_width=Pt(0.5))
txbox(s, "Key fix applied:  Per-horizon class weights  (1h: 13.8×  |  6h: 7.5×  |  24h: 2.9×)  +  predict() normalization",
      0.6, 5.95, 12.1, 0.35, size=13, color=ACCENT, bold=True)
txbox(s, "Without these fixes the model predicts 'normal' for everything and gets 99.4% accuracy — completely useless.",
      0.6, 6.26, 12.1, 0.3, size=12, color=LIGHT_GREY, italic=True)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "10 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Live Inference
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "Live Inference Results")
subheading(s, "Feeding real sensor windows to the trained model")
divider(s)

scenarios = [
    ("Normal rack (stable 22°C, healthy disk)",     "1h: 31.9%   6h: 41.8%   24h: 47.0%", False),
    ("Overheating  (temp 28°C → 55°C over 60 min)", "1h: 100.0%  6h:  99.9%  24h: 99.8%", True),
    ("Disk degrading  (health 0.85 → 0.40)",         "1h:  99.9%  6h:  99.9%  24h: 99.6%", True),
    ("Cooling failure  (temp 25°C → 48°C + power spike)", "1h: 100.0%  6h:  99.9%  24h: 99.9%", True),
    ("Power fluctuation  (5 kW → 12 kW spike)",     "1h:  67.4%  6h:  65.9%  24h: 57.4%", False),
]

for i, (name, probs, is_alert) in enumerate(scenarios):
    line_col = RGBColor(0xFF,0x44,0x44) if is_alert else ACCENT
    rect(s, 0.4, 1.3 + i*1.08, 12.5, 1.0, fill_color=CARD_BG, line_color=line_col, line_width=Pt(1))
    txbox(s, name, 0.6, 1.38 + i*1.08, 7.5, 0.35, size=14, bold=True,
          color=RGBColor(0xFF,0x66,0x66) if is_alert else WHITE)
    txbox(s, probs, 7.8, 1.42 + i*1.08, 4.9, 0.35, size=13,
          color=RGBColor(0xFF,0x44,0x44) if is_alert else NVIDIA_GREEN,
          bold=is_alert, align=PP_ALIGN.RIGHT)
    alert_txt = "  ← ALERT" if is_alert else ""
    if is_alert:
        txbox(s, "ALERT — threshold 70% exceeded", 0.6, 1.7 + i*1.08, 6, 0.28, size=11,
              color=RGBColor(0xFF,0x44,0x44), italic=True)
    else:
        txbox(s, "Normal — below 70% alert threshold", 0.6, 1.7 + i*1.08, 6, 0.28, size=11,
              color=LIGHT_GREY, italic=True)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "11 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Student Workflow
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "Student Workflow — How to Get Started")
subheading(s, "From zero to trained model in under 30 minutes")
divider(s)

steps = [
    ("1", "Get access from your instructor",
     "Ask your instructor to run:\n  bash deploy/instructor_grant_access.sh YOUR@gmail.com"),
    ("2", "Clone the repo",
     "git clone https://github.com/INSTRUCTOR/dc-world-model-tutorial\ncd dc-world-model-tutorial"),
    ("3", "Run the setup script",
     "GCS_BUCKET=instructor-bucket-name bash deploy/student_setup.sh\n→ installs gcloud, downloads 9.6 GB assets, installs Python deps"),
    ("4", "Generate synthetic failure data",
     "python3 deploy/06_generate_failure_data.py\n→ creates ~/sensor_timeseries.csv  (414,720 rows, ~25 MB)"),
    ("5", "Train the world model",
     "python3 deploy/07_world_model.py --csv ~/sensor_timeseries.csv --epochs 15\n→ trains on your GPU, saves model_output/best_model.pt"),
    ("6", "Run live inference",
     "See what failure probabilities your model predicts\nfor normal vs overheating vs disk degradation scenarios"),
]

for i, (num, title, detail) in enumerate(steps):
    col = i % 2
    row = i // 2
    l = 0.4 + col * 6.45
    t = 1.3 + row * 1.9
    rect(s, l, t, 6.15, 1.75, fill_color=CARD_BG, line_color=NVIDIA_GREEN, line_width=Pt(0.75))
    txbox(s, num, l+0.1, t+0.12, 0.45, 0.5, size=26, bold=True, color=NVIDIA_GREEN)
    txbox(s, title, l+0.55, t+0.1, 5.45, 0.35, size=14, bold=True, color=WHITE)
    txbox(s, detail, l+0.55, t+0.45, 5.45, 1.2, size=11, color=LIGHT_GREY)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "12 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Exercises
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "Exercises for Students")
subheading(s, "Three levels — work through them in order")
divider(s)

levels = [
    ("🟢  Beginner", NVIDIA_GREEN, [
        "Change FAILURE_RATE from 0.15 → 0.30 and retrain. Does accuracy improve?",
        "Add a 5th sensor feature: fan_rpm. Modify both the data generator and model.",
        "Change WINDOW_SIZE from 12 → 36 (3 hours of history). Does more context help?",
        "Print a confusion matrix for the 1h horizon predictions.",
    ]),
    ("🟡  Intermediate", ORANGE, [
        "Replace mean pooling with a [CLS] token like BERT. Does accuracy improve?",
        "Add early stopping: stop training if val_loss doesn't improve for 5 epochs.",
        "Write a script that sends a Slack alert when failure probability > 80%.",
        "Model cooling failure as correlated across an entire row of racks.",
    ]),
    ("🔴  Advanced", RGBColor(0xFF,0x44,0x44), [
        "Add a Transformer decoder that generates future sensor trajectories.",
        "Replace z-score normalization with global statistics saved in the checkpoint.",
        "Implement anomaly detection using reconstruction error (autoencoder).",
        "Federated learning: train on multiple data centers without sharing raw data.",
    ]),
]
for i, (title, color, items) in enumerate(levels):
    rect(s, 0.4 + i*4.3, 1.3, 4.1, 5.5, fill_color=CARD_BG, line_color=color, line_width=Pt(1))
    txbox(s, title, 0.55 + i*4.3, 1.38, 3.8, 0.4, size=15, bold=True, color=color)
    rect(s, 0.55 + i*4.3, 1.78, 3.8, 0.03, fill_color=color)
    for j, item in enumerate(items):
        rect(s, 0.55+i*4.3, 1.88+j*1.2, 3.8, 1.1, fill_color=RGBColor(0x0D,0x1B,0x2A),
             line_color=color, line_width=Pt(0.35))
        txbox(s, f"  {item}", 0.65+i*4.3, 1.93+j*1.2, 3.6, 1.0, size=12, color=LIGHT_GREY)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "13 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Key Concepts Glossary
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "Key Concepts Glossary")
divider(s, top=0.95)

terms = [
    ("USD", "Universal Scene Description — Pixar/NVIDIA's 3D file format. Scenes are assembled from many referenced sub-files."),
    ("Digital Twin", "A live virtual replica of a physical asset, synchronized with real sensor data and capable of simulation."),
    ("Omniverse Kit", "NVIDIA's platform for building apps that render USD scenes in real time and stream via WebRTC."),
    ("WebRTC", "Web protocol for real-time video streaming (used by Google Meet). We stream the 3D viewport to browsers."),
    ("GCS Fuse", "Mounts a GCS bucket as a local Linux folder — the VM reads USD from GCS as if it were on disk."),
    ("Transformer Encoder", "Neural network layer using attention to let every timestep relate to every other timestep."),
    ("World Model", "A neural network that predicts how a system will evolve — here, future failure probability per rack."),
    ("Vertex AI", "Google's managed ML platform: submit a training script, get a GPU, train, serve — no VM management."),
    ("Synthetic Data", "Data generated by simulation. Lets us create thousands of rare failure examples cheaply and safely."),
    ("Sliding Window", "Taking a fixed-length slice of a timeseries as model input — here, the last 60 minutes of sensor data."),
    ("Class Weights", "Penalty multiplier for under-represented classes. Prevents the model ignoring rare failure events."),
    ("Per-Window Norm", "Z-score normalization computed per sample window so the model learns shape/trend, not absolute value."),
]

for i, (term, defn) in enumerate(terms):
    col = i % 2
    row = i // 2
    t = 1.05 + row * 0.98
    l = 0.4 + col * 6.45
    rect(s, l, t, 6.15, 0.88, fill_color=CARD_BG, line_color=ACCENT, line_width=Pt(0.35))
    txbox(s, term, l+0.12, t+0.06, 2.0, 0.3, size=13, bold=True, color=NVIDIA_GREEN)
    txbox(s, defn, l+0.12, t+0.36, 5.9, 0.48, size=11, color=LIGHT_GREY)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "14 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Cost Summary
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)
heading(s, "Estimated Cloud Costs")
subheading(s, "Full tutorial from Phase 1 to Phase 9 — approximately $10 per student")
divider(s)

items = [
    ("GPU VM — g2-standard-8 + L4", "8 hours", "$0.40/hr", "$3.20"),
    ("GCS Storage — 10 GB/month",   "1 month", "$0.02/GB", "$0.20"),
    ("Artifact Registry — 10 GB",   "1 month", "$0.10/GB", "$0.10"),
    ("Vertex AI — A100 training",   "2 hours", "$3.00/hr", "$6.00"),
    ("Network egress (downloads)",  "~20 GB",  "$0.08/GB", "$1.60"),
]

headers = ["Resource", "Quantity", "Rate", "Cost"]
col_x = [0.5, 5.5, 8.5, 11.2]
col_w = [4.8, 2.8, 2.5, 1.8]

# Header row
rect(s, 0.4, 1.3, 12.5, 0.5, fill_color=NVIDIA_GREEN)
for j, (hdr, lx, lw) in enumerate(zip(headers, col_x, col_w)):
    txbox(s, hdr, lx, 1.35, lw, 0.4, size=14, bold=True, color=DARK_BG, align=PP_ALIGN.LEFT)

for i, row_data in enumerate(items):
    bg_col = CARD_BG if i % 2 == 0 else RGBColor(0x13,0x1D,0x35)
    rect(s, 0.4, 1.85 + i*0.65, 12.5, 0.6, fill_color=bg_col)
    for j, (cell, lx, lw) in enumerate(zip(row_data, col_x, col_w)):
        color = NVIDIA_GREEN if j == 3 else WHITE
        txbox(s, cell, lx, 1.9 + i*0.65, lw, 0.5, size=14, color=color,
              bold=(j==3), align=PP_ALIGN.LEFT)

# Total
rect(s, 0.4, 5.15, 12.5, 0.65, fill_color=RGBColor(0x1A,0x3A,0x1A), line_color=NVIDIA_GREEN, line_width=Pt(1.5))
txbox(s, "TOTAL (per student)", 0.6, 5.22, 5, 0.5, size=16, bold=True, color=WHITE)
txbox(s, "~$11.10", 10.5, 5.22, 2, 0.5, size=22, bold=True, color=NVIDIA_GREEN, align=PP_ALIGN.RIGHT)

# Cost tips
rect(s, 0.4, 5.95, 12.5, 0.95, fill_color=CARD_BG, line_color=ACCENT, line_width=Pt(0.5))
tips = "Cost tips:  Stop VM when not in use → gcloud compute instances stop datacenter-kit-vm --zone us-central1-a\nSet a $20 budget alert in GCP Billing console.  Students only need GCS access for assets (free tier covers most)."
txbox(s, tips, 0.6, 6.0, 12.1, 0.85, size=12, color=LIGHT_GREY)

txbox(s, "Atlanta Robotics — Digital Twin Tutorial", 0.5, 7.1, 9, 0.3, size=11, color=LIGHT_GREY, italic=True)
txbox(s, "15 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Closing / Thank You
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 0.08, fill_color=NVIDIA_GREEN)

txbox(s, "Build it.", 0.6, 1.4, 12, 1.0, size=64, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txbox(s, "Don't just study AI — ship it.",
      0.6, 2.4, 12, 0.6, size=28, color=ACCENT, align=PP_ALIGN.LEFT)

rect(s, 0.6, 3.1, 11.5, 0.05, fill_color=NVIDIA_GREEN)

links = [
    ("GitHub Repo",    "github.com/YOUR_USERNAME/dc-world-model-tutorial"),
    ("NVIDIA Omniverse", "developer.nvidia.com/omniverse"),
    ("GCP Console",    "console.cloud.google.com"),
    ("PyTorch Docs",   "pytorch.org/docs"),
]
for i, (label, url) in enumerate(links):
    col = i % 2
    row = i // 2
    txbox(s, f"{label}:", 0.6 + col*6.2, 3.3 + row*0.55, 2.2, 0.45, size=14, bold=True, color=NVIDIA_GREEN)
    txbox(s, url, 2.7 + col*6.2, 3.3 + row*0.55, 3.8, 0.45, size=14, color=LIGHT_GREY)

rect(s, 0.6, 4.6, 11.5, 0.05, fill_color=NVIDIA_GREEN)

txbox(s, "Atlanta Robotics", 0.6, 4.8, 6, 0.6, size=28, bold=True, color=WHITE)
txbox(s, "College Student Tutorial Series — 2026",
      0.6, 5.35, 8, 0.45, size=18, color=LIGHT_GREY)

txbox(s, "Questions? Open a GitHub Issue on the repo.",
      0.6, 6.0, 10, 0.4, size=16, color=ACCENT, italic=True)

rect(s, 0, 7.42, 13.33, 0.08, fill_color=NVIDIA_GREEN)
txbox(s, "16 / 16", 12, 7.1, 1, 0.3, size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
